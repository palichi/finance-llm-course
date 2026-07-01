"""PPO + Advanced RAG 트레이딩 시스템 설계 PPT 생성."""
from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
from pptx.util import Inches, Pt
import copy

# ── 색상 팔레트 ──────────────────────────────────────────────────────────────
C_BG        = RGBColor(0x0D, 0x1B, 0x2A)   # 네이비 배경
C_ACCENT    = RGBColor(0x00, 0xB4, 0xD8)   # 밝은 파랑 (포인트)
C_ACCENT2   = RGBColor(0x90, 0xE0, 0xEF)   # 연한 파랑
C_GOLD      = RGBColor(0xFF, 0xD1, 0x66)   # 골드
C_GREEN     = RGBColor(0x52, 0xB7, 0x88)   # 그린
C_WHITE     = RGBColor(0xFF, 0xFF, 0xFF)
C_GRAY      = RGBColor(0xB0, 0xBE, 0xC5)
C_RED       = RGBColor(0xFF, 0x6B, 0x6B)
C_DARK_BOX  = RGBColor(0x1A, 0x2E, 0x44)   # 카드 배경


def px(n): return Pt(n)
def inch(n): return Inches(n)


def set_bg(slide, prs, color: RGBColor):
    """슬라이드 배경색 설정."""
    from pptx.oxml.ns import qn
    from lxml import etree
    bg = slide.background
    fill = bg.fill
    fill.solid()
    fill.fore_color.rgb = color


def add_rect(slide, left, top, w, h, fill_color, alpha=None):
    """색상 박스 추가."""
    shape = slide.shapes.add_shape(
        1,  # MSO_SHAPE_TYPE.RECTANGLE
        inch(left), inch(top), inch(w), inch(h)
    )
    shape.line.fill.background()
    shape.fill.solid()
    shape.fill.fore_color.rgb = fill_color
    return shape


def add_text(slide, text, left, top, w, h,
             font_size=18, bold=False, color=C_WHITE,
             align=PP_ALIGN.LEFT, wrap=True):
    """텍스트 박스 추가."""
    txBox = slide.shapes.add_textbox(inch(left), inch(top), inch(w), inch(h))
    tf = txBox.text_frame
    tf.word_wrap = wrap
    p = tf.paragraphs[0]
    p.alignment = align
    run = p.add_run()
    run.text = text
    run.font.size = px(font_size)
    run.font.bold = bold
    run.font.color.rgb = color
    run.font.name = "Malgun Gothic"
    return txBox


def add_multiline(slide, lines, left, top, w, h,
                  font_size=14, color=C_WHITE, line_color=None, bold_first=False):
    """여러 줄 텍스트 박스."""
    txBox = slide.shapes.add_textbox(inch(left), inch(top), inch(w), inch(h))
    tf = txBox.text_frame
    tf.word_wrap = True
    for i, (line, fc, sz, bd) in enumerate(lines):
        if i == 0:
            p = tf.paragraphs[0]
        else:
            p = tf.add_paragraph()
        p.alignment = PP_ALIGN.LEFT
        run = p.add_run()
        run.text = line
        run.font.size = px(sz)
        run.font.bold = bd
        run.font.color.rgb = fc
        run.font.name = "Malgun Gothic"


def slide_title_only(prs, title, subtitle=None):
    """타이틀 슬라이드."""
    slide = prs.slides.add_slide(prs.slide_layouts[6])  # blank
    set_bg(slide, prs, C_BG)

    # 상단 강조 바
    add_rect(slide, 0, 0, 10, 0.08, C_ACCENT)

    # 중앙 타이틀 박스
    add_rect(slide, 0.5, 2.0, 9.0, 1.2, C_DARK_BOX)
    add_text(slide, title, 0.6, 2.1, 8.8, 1.0,
             font_size=36, bold=True, color=C_ACCENT, align=PP_ALIGN.CENTER)

    if subtitle:
        add_text(slide, subtitle, 0.5, 3.4, 9.0, 0.6,
                 font_size=18, color=C_ACCENT2, align=PP_ALIGN.CENTER)

    # 하단 바
    add_rect(slide, 0, 7.12, 10, 0.08, C_ACCENT)
    add_text(slide, "PPO 강화학습 + Advanced RAG  |  금융 LLM 실습 Day 3",
             0, 7.2, 10, 0.3, font_size=11, color=C_GRAY, align=PP_ALIGN.CENTER)
    return slide


def slide_section(prs, num, title, color=C_ACCENT):
    """섹션 구분 슬라이드."""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_bg(slide, prs, C_BG)
    add_rect(slide, 0, 0, 10, 0.08, color)

    add_rect(slide, 1.5, 2.5, 7.0, 2.2, C_DARK_BOX)
    add_text(slide, f"STEP {num}", 1.5, 2.6, 7.0, 0.6,
             font_size=20, bold=False, color=color, align=PP_ALIGN.CENTER)
    add_text(slide, title, 1.5, 3.1, 7.0, 1.2,
             font_size=30, bold=True, color=C_WHITE, align=PP_ALIGN.CENTER)

    add_rect(slide, 0, 7.12, 10, 0.08, color)
    return slide


def slide_content(prs, title, content_func):
    """일반 콘텐츠 슬라이드 — content_func(slide) 콜백으로 내용 채움."""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_bg(slide, prs, C_BG)
    add_rect(slide, 0, 0, 10, 0.08, C_ACCENT)
    add_rect(slide, 0, 0.08, 10, 0.72, C_DARK_BOX)
    add_text(slide, title, 0.3, 0.1, 9.4, 0.68,
             font_size=22, bold=True, color=C_ACCENT)
    add_rect(slide, 0, 7.12, 10, 0.08, C_ACCENT)
    add_text(slide, "PPO 강화학습 + Advanced RAG  |  금융 LLM 실습 Day 3",
             0, 7.2, 10, 0.3, font_size=11, color=C_GRAY, align=PP_ALIGN.CENTER)
    content_func(slide)
    return slide


# ============================================================
# 슬라이드별 콘텐츠 정의
# ============================================================

def make_ppt():
    prs = Presentation()
    prs.slide_width  = Inches(10)
    prs.slide_height = Inches(7.5)

    # ── 0. 커버 ──────────────────────────────────────────────────────────────
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_bg(slide, prs, C_BG)
    add_rect(slide, 0, 0, 10, 0.1, C_ACCENT)
    add_rect(slide, 0, 7.4, 10, 0.1, C_ACCENT)

    # 좌측 세로 바
    add_rect(slide, 0, 0, 0.12, 7.5, C_ACCENT)

    add_text(slide, "PPO 강화학습", 0.5, 1.4, 9, 0.9,
             font_size=42, bold=True, color=C_WHITE)
    add_text(slide, "+  Advanced RAG", 0.5, 2.2, 9, 0.8,
             font_size=36, bold=True, color=C_ACCENT)
    add_text(slide, "주식 매매 의사결정 시스템 설계 전 과정",
             0.5, 3.1, 9, 0.6, font_size=20, color=C_ACCENT2)

    add_rect(slide, 0.5, 3.9, 5.5, 0.04, C_ACCENT)

    add_text(slide, "데이터 수집 → 지표 계산 → 환경 구축 → 학습 → 추론 → RAG → LLM 설명 → 챗봇",
             0.5, 4.05, 9, 0.5, font_size=13, color=C_GRAY)
    add_text(slide, "금융 LLM 실습  Day 3  |  01_golden_cross",
             0.5, 6.7, 9, 0.4, font_size=13, color=C_GRAY)

    # ── 1. 전체 아키텍처 ──────────────────────────────────────────────────────
    def s1(slide):
        # 주 트랙 박스
        add_rect(slide, 0.3, 0.9, 4.2, 5.7, RGBColor(0x0A, 0x25, 0x3D))
        add_text(slide, "주(主)  PPO 트랙", 0.4, 0.95, 4.0, 0.45,
                 font_size=15, bold=True, color=C_GOLD)

        main_steps = [
            ("① 데이터 수집", "stock_prices.csv\n200종목 · 2년치"),
            ("② 기술 지표",   "SMA / RSI / 이격도\n골든·데드크로스"),
            ("③ 거래 환경",   "gymnasium TradingEnv\n종목별 episode 분리"),
            ("④ PPO 학습",   "stable-baselines3\n단일 정책망 200종목"),
            ("⑤ 추론",       "매수 / 유보 / 매도\n+ 확률 분포"),
        ]
        for i, (t, d) in enumerate(main_steps):
            y = 1.45 + i * 1.0
            add_rect(slide, 0.4, y, 4.0, 0.8, C_DARK_BOX)
            add_text(slide, t, 0.5, y + 0.02, 3.8, 0.35,
                     font_size=13, bold=True, color=C_ACCENT)
            add_text(slide, d, 0.5, y + 0.38, 3.8, 0.4,
                     font_size=11, color=C_GRAY)
            if i < 4:
                add_text(slide, "↓", 2.3, y + 0.82, 0.4, 0.2,
                         font_size=14, color=C_ACCENT, align=PP_ALIGN.CENTER)

        # 부 트랙 박스
        add_rect(slide, 5.1, 0.9, 4.5, 5.7, RGBColor(0x0A, 0x25, 0x3D))
        add_text(slide, "부(副)  Advanced RAG 트랙", 5.2, 0.95, 4.3, 0.45,
                 font_size=15, bold=True, color=C_GREEN)

        rag_steps = [
            ("⑥ 코퍼스 구축",  "사례 카드 생성\nPPO 판단 + 사후 수익률"),
            ("⑦ 임베딩 저장",  "KR-SBERT → ChromaDB\n+ BM25 인덱스"),
            ("⑧ Advanced RAG", "Query Expansion\nHybrid Search · RRF · Rerank"),
            ("⑨ LLM 설명",    "Claude Haiku\n자연어 한국어 설명"),
            ("⑩ 챗봇 앱",     "Streamlit\n종목명 입력 → 판단 + 근거"),
        ]
        for i, (t, d) in enumerate(rag_steps):
            y = 1.45 + i * 1.0
            add_rect(slide, 5.2, y, 4.3, 0.8, C_DARK_BOX)
            add_text(slide, t, 5.3, y + 0.02, 4.1, 0.35,
                     font_size=13, bold=True, color=C_GREEN)
            add_text(slide, d, 5.3, y + 0.38, 4.1, 0.4,
                     font_size=11, color=C_GRAY)
            if i < 4:
                add_text(slide, "↓", 7.3, y + 0.82, 0.4, 0.2,
                         font_size=14, color=C_GREEN, align=PP_ALIGN.CENTER)

        # 중앙 연결 화살표
        add_text(slide, "→", 4.5, 5.0, 0.6, 0.5,
                 font_size=24, bold=True, color=C_GOLD, align=PP_ALIGN.CENTER)
        add_text(slide, "PPO 결정\n후 호출", 4.37, 5.5, 0.86, 0.4,
                 font_size=9, color=C_GOLD, align=PP_ALIGN.CENTER)

    slide_content(prs, "전체 시스템 아키텍처 — PPO가 두뇌, RAG는 설명", s1)

    # ── 2. 데이터 수집 ────────────────────────────────────────────────────────
    def s2(slide):
        add_text(slide, "금융위원회 공공데이터 포털 — stock_prices.csv",
                 0.3, 0.9, 9.4, 0.4, font_size=14, color=C_ACCENT2)

        cols = [
            ("srtnCd",   "종목코드(6자리)",   "종목 구분 기준 키"),
            ("basDt",    "기준일자 YYYYMMDD", "시간축·정렬 기준"),
            ("clpr",     "종가",              "지표 계산의 핵심값"),
            ("mkp/hipr/lopr", "시가·고가·저가", "State 구성 참고"),
            ("trqu",     "거래량",            "신호 품질 향상"),
            ("itmsNm",   "종목명",            "챗봇 질의 매핑"),
        ]

        header_y = 1.35
        add_rect(slide, 0.3, header_y, 2.5, 0.38, C_ACCENT)
        add_rect(slide, 2.8, header_y, 3.2, 0.38, C_ACCENT)
        add_rect(slide, 6.0, header_y, 3.7, 0.38, C_ACCENT)
        add_text(slide, "컬럼명",    0.3, header_y + 0.05, 2.5, 0.3,
                 font_size=13, bold=True, color=C_BG, align=PP_ALIGN.CENTER)
        add_text(slide, "의미",      2.8, header_y + 0.05, 3.2, 0.3,
                 font_size=13, bold=True, color=C_BG, align=PP_ALIGN.CENTER)
        add_text(slide, "실습에서의 역할", 6.0, header_y + 0.05, 3.7, 0.3,
                 font_size=13, bold=True, color=C_BG, align=PP_ALIGN.CENTER)

        for i, (c, m, r) in enumerate(cols):
            y = 1.78 + i * 0.47
            bg = C_DARK_BOX if i % 2 == 0 else RGBColor(0x16, 0x28, 0x3C)
            add_rect(slide, 0.3, y, 2.5, 0.42, bg)
            add_rect(slide, 2.8, y, 3.2, 0.42, bg)
            add_rect(slide, 6.0, y, 3.7, 0.42, bg)
            add_text(slide, c, 0.4, y + 0.07, 2.3, 0.3,
                     font_size=12, bold=True, color=C_ACCENT)
            add_text(slide, m, 2.9, y + 0.07, 3.0, 0.3, font_size=12, color=C_WHITE)
            add_text(slide, r, 6.1, y + 0.07, 3.5, 0.3, font_size=12, color=C_GRAY)

        # 통계 박스
        add_rect(slide, 0.3, 4.7, 9.4, 0.65, C_DARK_BOX)
        add_text(slide,
                 "종목 수: 200개  ·  기간: 2024-01-02 ~ 2026-06-18  ·  "
                 "종목당 행 수: 138 ~ 598  ·  전체 약 117,000행",
                 0.5, 4.78, 9.0, 0.5,
                 font_size=13, color=C_GOLD, align=PP_ALIGN.CENTER)

        add_text(slide,
                 "※ Long format — 한 파일에 전 종목, srtnCd 컬럼으로 종목 구분 (groupby 필수)",
                 0.3, 5.5, 9.4, 0.35, font_size=12, color=C_RED)

    slide_content(prs, "STEP 1  데이터 수집 · 구조 파악", s2)

    # ── 3. 기술 지표 ──────────────────────────────────────────────────────────
    def s3(slide):
        indicators = [
            ("SMA 5/20/60",      "단순 이동평균",
             "골든·데드크로스 기준\nState 이격률 계산"),
            ("EMA 20",           "지수 이동평균",
             "최근 가격에 더 높은 가중치"),
            ("골든·데드크로스", "SMA5 vs SMA20 돌파",
             "golden_flag / dead_flag (0/1)"),
            ("이격도20",         "(종가/SMA20)×100",
             "과열≥105, 침체≤95, 중립"),
            ("RSI 14",           "Wilder's Smoothing",
             "과매수≥70, 과매도≤30\nrsi14_norm = (RSI-50)/50"),
        ]
        for i, (name, formula, note) in enumerate(indicators):
            col = i % 3
            row = i // 3
            x = 0.3 + col * 3.25
            y = 1.0 + row * 2.7
            add_rect(slide, x, y, 3.1, 2.4, C_DARK_BOX)
            add_rect(slide, x, y, 3.1, 0.38, C_ACCENT)
            add_text(slide, name, x + 0.1, y + 0.04, 2.9, 0.32,
                     font_size=14, bold=True, color=C_BG)
            add_text(slide, formula, x + 0.1, y + 0.45, 2.9, 0.45,
                     font_size=12, color=C_GOLD)
            add_text(slide, note, x + 0.1, y + 0.95, 2.9, 1.2,
                     font_size=11, color=C_GRAY)

        # 핵심 원칙
        add_rect(slide, 0.3, 5.65, 9.4, 0.7, RGBColor(0x00, 0x4E, 0x6E))
        add_text(slide, "핵심 원칙: srtnCd 별 groupby 계산 — 종목 경계를 절대 넘지 않는다",
                 0.5, 5.72, 9.0, 0.3, font_size=13, bold=True, color=C_WHITE)
        add_text(slide, "신규상장 등 lookback 20일 미만 종목/구간 → 학습·추론 대상에서 제외",
                 0.5, 6.05, 9.0, 0.3, font_size=12, color=C_ACCENT2)

    slide_content(prs, "STEP 2  기술 지표 계산 (indicators/technical.py)", s3)

    # ── 4. State / Action / Reward ────────────────────────────────────────────
    def s4(slide):
        # State
        add_rect(slide, 0.3, 0.95, 3.1, 5.7, C_DARK_BOX)
        add_rect(slide, 0.3, 0.95, 3.1, 0.42, C_ACCENT)
        add_text(slide, "State  (9차원 × 20일)", 0.4, 0.98, 2.9, 0.35,
                 font_size=13, bold=True, color=C_BG)
        state_items = [
            "SMA5 이격률",
            "SMA20 이격률",
            "SMA60 이격률",
            "골든크로스 플래그",
            "데드크로스 플래그",
            "이격도20 (중심화)",
            "RSI14 (정규화)",
            "현재 포지션 (0/1)",
            "미실현 손익",
        ]
        for i, s in enumerate(state_items):
            add_text(slide, f"• {s}", 0.4, 1.44 + i * 0.47, 2.9, 0.42,
                     font_size=12, color=C_WHITE)

        # Action
        add_rect(slide, 3.6, 0.95, 2.8, 5.7, C_DARK_BOX)
        add_rect(slide, 3.6, 0.95, 2.8, 0.42, C_GREEN)
        add_text(slide, "Action  (Discrete 3)", 3.7, 0.98, 2.6, 0.35,
                 font_size=13, bold=True, color=C_BG)
        actions = [
            ("0", "매도", C_RED),
            ("1", "유보", C_GOLD),
            ("2", "매수", C_GREEN),
        ]
        for i, (n, label, c) in enumerate(actions):
            y = 1.55 + i * 1.4
            add_rect(slide, 3.75, y, 2.5, 0.95, RGBColor(0x0A, 0x20, 0x32))
            add_text(slide, n, 3.85, y + 0.08, 0.5, 0.5,
                     font_size=28, bold=True, color=c)
            add_text(slide, label, 4.4, y + 0.25, 1.8, 0.45,
                     font_size=20, bold=True, color=c)

        # Reward
        add_rect(slide, 6.6, 0.95, 3.1, 5.7, C_DARK_BOX)
        add_rect(slide, 6.6, 0.95, 3.1, 0.42, C_GOLD)
        add_text(slide, "Reward 설계", 6.7, 0.98, 2.9, 0.35,
                 font_size=13, bold=True, color=C_BG)

        reward_lines = [
            ("기본 보상", C_WHITE, 13, True),
            ("자산가치 변화율", C_ACCENT2, 12, False),
            ("− 거래비용 0.2~0.3%", C_ACCENT2, 12, False),
            ("", C_WHITE, 8, False),
            ("셰이핑 (Shaping)", C_GOLD, 13, True),
            ("RSI>70 매수 시  −0.10", C_RED, 12, False),
            ("골든크로스 보유 유지  +0.05", C_GREEN, 12, False),
            ("이격도>110 매수 시  −0.10", C_RED, 12, False),
            ("", C_WHITE, 8, False),
            ("가중치 → config 분리", C_GRAY, 11, False),
            ("on/off 플래그 제공", C_GRAY, 11, False),
            ("발동 횟수 카운트 로깅", C_GRAY, 11, False),
        ]
        add_multiline(slide, [(t, c, s, b) for t, c, s, b in reward_lines],
                      6.7, 1.45, 2.9, 5.0)

    slide_content(prs, "STEP 3  강화학습 설계 — State · Action · Reward", s4)

    # ── 5. 학습 전략 ──────────────────────────────────────────────────────────
    def s5(slide):
        # 모델 1개로 200종목
        add_rect(slide, 0.3, 0.95, 9.4, 0.95, C_DARK_BOX)
        add_text(slide, "모델 1개 — 200종목의 Episode를 섞어서 단일 정책망 학습",
                 0.5, 1.0, 9.0, 0.45, font_size=15, bold=True, color=C_ACCENT)
        add_text(slide,
                 "종목코드(srtnCd)는 State에 포함하지 않음 → 지표 패턴의 일반 규칙 학습",
                 0.5, 1.45, 9.0, 0.38, font_size=13, color=C_GRAY)

        # 데이터 분할
        add_rect(slide, 0.3, 2.05, 4.5, 2.3, C_DARK_BOX)
        add_text(slide, "데이터 분할 전략", 0.4, 2.1, 4.3, 0.38,
                 font_size=14, bold=True, color=C_GOLD)
        split_lines = [
            ("Train",       "~2025-12-31 (약 80%)",  C_GREEN),
            ("Validation",  "2026-01-01~ (약 20%)",  C_ACCENT),
            ("룩어헤드 방지", "종목 내 시간 기준 분리",  C_GRAY),
            ("증분 평가",    "N스텝마다 val 성과 측정", C_GRAY),
        ]
        for i, (k, v, c) in enumerate(split_lines):
            add_text(slide, k, 0.4, 2.56 + i * 0.42, 1.6, 0.38,
                     font_size=12, bold=True, color=c)
            add_text(slide, v, 2.1, 2.56 + i * 0.42, 2.6, 0.38,
                     font_size=12, color=C_WHITE)

        # PPO 하이퍼파라미터
        add_rect(slide, 5.0, 2.05, 4.7, 2.3, C_DARK_BOX)
        add_text(slide, "PPO 하이퍼파라미터 (config 분리)", 5.1, 2.1, 4.5, 0.38,
                 font_size=14, bold=True, color=C_GOLD)
        params = [
            ("알고리즘",  "PPO (stable-baselines3)"),
            ("정책망",   "MlpPolicy  [64, 64]"),
            ("n_steps",  "2048  · batch_size 64"),
            ("gamma",    "0.99  · lr 3e-4"),
        ]
        for i, (k, v) in enumerate(params):
            add_text(slide, k, 5.1, 2.56 + i * 0.42, 1.6, 0.38,
                     font_size=12, bold=True, color=C_ACCENT2)
            add_text(slide, v, 6.85, 2.56 + i * 0.42, 2.75, 0.38,
                     font_size=12, color=C_WHITE)

        # 평가 지표
        add_rect(slide, 0.3, 4.5, 9.4, 1.95, C_DARK_BOX)
        add_text(slide, "Validation 성과 평가 지표", 0.5, 4.56, 4.0, 0.38,
                 font_size=14, bold=True, color=C_GREEN)
        metrics = ["누적 수익률", "MDD (최대 낙폭)", "샤프 비율", "승률"]
        for i, m in enumerate(metrics):
            x = 0.5 + i * 2.35
            add_rect(slide, x, 4.97, 2.1, 0.9, RGBColor(0x0A, 0x20, 0x32))
            add_text(slide, m, x + 0.05, 5.1, 2.0, 0.65,
                     font_size=13, color=C_ACCENT2, align=PP_ALIGN.CENTER)
        add_text(slide, "최고 Validation 성과 체크포인트 → models/best_model.zip 저장",
                 0.5, 6.05, 9.0, 0.35, font_size=12, color=C_GOLD)

    slide_content(prs, "STEP 4-5  PPO 학습 전략 (train/run_training.py)", s5)

    # ── 6. 추론 + 수치 근거 ──────────────────────────────────────────────────
    def s6(slide):
        # 파이프라인 흐름
        steps = [
            ("종목코드 / 종목명 입력", C_ACCENT),
            ("최근 20일 데이터 로드", C_ACCENT),
            ("지표 계산 → State 구성", C_ACCENT),
            ("PPO 모델 추론", C_GOLD),
            ("Action + 확률 분포 반환", C_GOLD),
        ]
        for i, (label, c) in enumerate(steps):
            x = 0.25 + i * 1.9
            add_rect(slide, x, 1.0, 1.7, 0.7, C_DARK_BOX)
            add_text(slide, label, x + 0.05, 1.08, 1.6, 0.54,
                     font_size=11, color=c, align=PP_ALIGN.CENTER)
            if i < 4:
                add_text(slide, "→", x + 1.72, 1.2, 0.2, 0.35,
                         font_size=14, bold=True, color=C_ACCENT, align=PP_ALIGN.CENTER)

        # rule_based 출력 항목
        add_rect(slide, 0.3, 1.95, 4.5, 4.3, C_DARK_BOX)
        add_text(slide, "explain/rule_based.py 출력 (LLM 미사용 · 결정론적)",
                 0.4, 2.0, 4.3, 0.38, font_size=13, bold=True, color=C_ACCENT2)
        rb_items = [
            ("RSI14 구간",    "과매수(≥70) / 과매도(≤30) / 중립",     C_WHITE),
            ("이격도20 구간",  "과열(≥105) / 침체(≤95) / 중립",       C_WHITE),
            ("크로스 여부",    "golden_flag / dead_flag",              C_WHITE),
            ("Shaping 발동",  "RSI 매수 페널티 / 골든크로스 보너스",   C_GOLD),
            ("확신도 낮음",   "1·2위 확률 차 < 10%p",                 C_RED),
            ("action_probs",  "BUY / HOLD / SELL 확률 전체",           C_ACCENT2),
        ]
        for i, (k, v, c) in enumerate(rb_items):
            y = 2.45 + i * 0.6
            add_text(slide, k, 0.4, y, 1.6, 0.52,
                     font_size=12, bold=True, color=c)
            add_text(slide, v, 2.1, y, 2.6, 0.52,
                     font_size=11, color=C_GRAY)

        # 출력 예시
        add_rect(slide, 5.0, 1.95, 4.7, 4.3, RGBColor(0x06, 0x14, 0x22))
        add_text(slide, "출력 예시", 5.1, 2.0, 4.5, 0.38,
                 font_size=13, bold=True, color=C_GREEN)
        example = (
            "[000080  하이트진로]\n"
            "기준일: 2025-06-05\n"
            "최근 5일 종가: 22,300 → 22,150\n\n"
            "판단: 매수  (확률 62%)\n"
            "유보 31%  ·  매도 7%\n\n"
            "RSI14 = 58.2  (중립)\n"
            "이격도20 = 102.1  (중립)\n"
            "골든크로스 4일 전 발생\n"
            "확신도: 보통 (차이 31%p)"
        )
        add_text(slide, example, 5.1, 2.45, 4.5, 3.7,
                 font_size=12, color=C_ACCENT2)

    slide_content(prs, "STEP 6  추론 + 수치 근거 (inference · explain/rule_based)", s6)

    # ── 7. 왜 Advanced RAG인가 ────────────────────────────────────────────────
    def s7(slide):
        add_text(slide,
                 "PPO는 '매수'를 결정하지만, '과거에 이런 패턴이 있었는가? 그때 어떻게 됐는가?'는 알지 못합니다.",
                 0.3, 0.92, 9.4, 0.45, font_size=13, color=C_ACCENT2)

        problems = [
            ("수치 조건 검색 약함",
             "벡터 검색만으로는 'RSI=70'처럼\n정확한 숫자 조건을 놓칠 수 있다",
             "Hybrid Search\n(벡터 + BM25 키워드)"),
            ("검색 커버리지 좁음",
             "'지금 사도 될까?' 한 문장만으론\n관련 사례를 다 못 찾음",
             "Query Expansion\n(조건별 쿼리 5개로 확장)"),
            ("관련 없는 사례 혼입",
             "의미적으로만 비슷한 게 아니라\n수치 조건도 맞아야 함",
             "Metadata Filtering\n(RSI/이격도 범위 · 날짜 필터)"),
            ("1차 검색 노이즈",
             "벡터 유사도 1위가 실제로는\n관련성이 낮을 수 있음",
             "Reranking\n(코사인 유사도 재정렬)"),
        ]
        for i, (prob, detail, sol) in enumerate(problems):
            col = i % 2
            row = i // 2
            x = 0.3 + col * 4.85
            y = 1.5 + row * 2.55

            # 문제 박스
            add_rect(slide, x, y, 4.55, 2.3, C_DARK_BOX)
            add_rect(slide, x, y, 4.55, 0.38, RGBColor(0x6B, 0x1F, 0x1F))
            add_text(slide, f"문제: {prob}", x + 0.1, y + 0.04, 4.3, 0.32,
                     font_size=12, bold=True, color=C_WHITE)
            add_text(slide, detail, x + 0.1, y + 0.45, 4.35, 0.75,
                     font_size=11, color=C_GRAY)

            # 해법 화살표
            add_text(slide, "→ 해법:", x + 0.1, y + 1.22, 0.9, 0.3,
                     font_size=11, color=C_GOLD)
            add_text(slide, sol, x + 1.1, y + 1.22, 3.35, 0.8,
                     font_size=12, bold=True, color=C_GREEN)

    slide_content(prs, "STEP 6-1  왜 Advanced RAG가 필요한가", s7)

    # ── 8. Corpus 구축 ────────────────────────────────────────────────────────
    def s8(slide):
        # 사례 카드 예시
        add_rect(slide, 0.3, 0.92, 9.4, 1.35, RGBColor(0x06, 0x18, 0x2A))
        add_text(slide, "사례 카드 텍스트 형식", 0.5, 0.95, 4.0, 0.38,
                 font_size=13, bold=True, color=C_GOLD)
        card_text = (
            "하이트진로(000080) 2025-06-05: RSI14=58.2(중립), 이격도20=102.1(중립), "
            "골든크로스=있음, 데드크로스=없음, 당시 PPO 판단=매수(확률 62%), "
            "관련 공시 없음.  이후 5일 수익률=+2.1%, 10일 수익률=+3.8%"
        )
        add_text(slide, card_text, 0.4, 1.35, 9.3, 0.85,
                 font_size=11, color=C_ACCENT2)

        # 구축 흐름
        flow_items = [
            ("주가 CSV\n로드", C_ACCENT),
            ("지표 계산\n(종목별)", C_ACCENT),
            ("PPO 추론\n(각 날짜)", C_GOLD),
            ("사례 카드\n텍스트 생성", C_GOLD),
            ("KR-SBERT\n임베딩", C_GREEN),
            ("ChromaDB\n+ BM25 저장", C_GREEN),
        ]
        for i, (label, c) in enumerate(flow_items):
            x = 0.25 + i * 1.6
            add_rect(slide, x, 2.5, 1.45, 0.95, C_DARK_BOX)
            add_text(slide, label, x + 0.05, 2.6, 1.35, 0.75,
                     font_size=11, color=c, align=PP_ALIGN.CENTER)
            if i < 5:
                add_text(slide, "→", x + 1.47, 2.85, 0.15, 0.3,
                         font_size=12, bold=True, color=C_ACCENT, align=PP_ALIGN.CENTER)

        # 메타데이터
        add_rect(slide, 0.3, 3.65, 4.5, 2.65, C_DARK_BOX)
        add_text(slide, "ChromaDB 저장 메타데이터", 0.4, 3.7, 4.3, 0.38,
                 font_size=13, bold=True, color=C_ACCENT2)
        meta_items = [
            "ticker · name · date",
            "rsi14 · disparity20",
            "golden_flag · dead_flag",
            "ppo_action · ppo_prob",
            "ret5 · ret10 (사후 수익률)",
        ]
        for i, m in enumerate(meta_items):
            add_text(slide, f"• {m}", 0.5, 4.15 + i * 0.44, 4.1, 0.4,
                     font_size=12, color=C_WHITE)

        # BM25 박스
        add_rect(slide, 5.0, 3.65, 4.7, 2.65, C_DARK_BOX)
        add_text(slide, "BM25 키워드 인덱스", 5.1, 3.7, 4.5, 0.38,
                 font_size=13, bold=True, color=C_ACCENT2)
        bm_items = [
            "rank_bm25 (BM25Okapi)",
            "한/영 숫자 토크나이징",
            "corpus/bm25_index.pkl",
            "증분 업데이트 지원",
            "Hybrid Search에서 병렬 사용",
        ]
        for i, m in enumerate(bm_items):
            add_text(slide, f"• {m}", 5.1, 4.15 + i * 0.44, 4.5, 0.4,
                     font_size=12, color=C_WHITE)

        add_text(slide,
                 "사후 수익률(5/10일)은 '기록' 목적 — PPO 학습에 미사용 (데이터 누설 아님)",
                 0.3, 6.38, 9.4, 0.35, font_size=12, color=C_RED)

    slide_content(prs, "STEP 7-A  RAG Corpus 구축 (corpus/build_corpus.py)", s8)

    # ── 9. Advanced RAG 4단계 ────────────────────────────────────────────────
    def s9(slide):
        steps_rag = [
            ("1  Query Expansion",
             "규칙 기반 템플릿으로 쿼리 최대 5개 생성\n"
             "예) 'RSI 과매수 매수 사례' · '골든크로스 직후 유보 사례'\n"
             "LLM 대신 규칙 템플릿 → 속도·비용 우선",
             C_ACCENT),
            ("2  Hybrid Search",
             "벡터 검색 (ChromaDB cosine)\n+ BM25 키워드 검색 동시 수행\n"
             "RRF (Reciprocal Rank Fusion)로 점수 결합\n"
             "스케일 정규화 불필요 · 이상치에 강건",
             C_GREEN),
            ("3  Metadata Filtering",
             "RSI ±10 · 이격도 ±5 범위 내 사례만 통과\n"
             "자기참조 방지: 현재 날짜 이후 데이터 제외\n"
             "후보 3건 미만 → 범위 2배/4배 단계적 완화\n"
             "(Fallback 3단계)",
             C_GOLD),
            ("4  Reranking",
             "쿼리 임베딩 vs 사례 임베딩 코사인 유사도\n"
             "ChromaDB 임베딩 재사용 → 추가 forward 없음\n"
             "최종 점수: 0.6×코사인 + 0.4×RRF×100\n"
             "Top-K (기본 5건) 반환",
             C_RED),
        ]
        for i, (title, detail, c) in enumerate(steps_rag):
            col = i % 2
            row = i // 2
            x = 0.3 + col * 4.85
            y = 0.95 + row * 2.9

            add_rect(slide, x, y, 4.55, 2.65, C_DARK_BOX)
            add_rect(slide, x, y, 4.55, 0.45, RGBColor(0x0A, 0x20, 0x32))
            add_rect(slide, x, y, 0.07, 2.65, c)
            add_text(slide, title, x + 0.15, y + 0.06, 4.3, 0.35,
                     font_size=14, bold=True, color=c)
            add_text(slide, detail, x + 0.15, y + 0.52, 4.3, 2.05,
                     font_size=11, color=C_WHITE)

        add_text(slide, "모든 모델·인덱스 — 모듈 로드 시 1회만 초기화 (매 호출마다 재로딩 금지)",
                 0.3, 6.9, 9.4, 0.35, font_size=12, color=C_GRAY,
                 align=PP_ALIGN.CENTER)

    slide_content(prs, "STEP 7-B  Advanced RAG 4단계 파이프라인 (rag_retriever.py)", s9)

    # ── 10. LLM 설명 생성 ────────────────────────────────────────────────────
    def s10(slide):
        # 흐름
        pipeline = [
            ("ExplainResult\n(수치 근거)", C_ACCENT),
            ("RAGResult\n(Top-5 사례)", C_GREEN),
            ("프롬프트\n조립", C_GOLD),
            ("Claude\nHaiku 4-5", C_RED),
            ("검증\n(숫자·인용)", C_GOLD),
            ("자연어\n설명 출력", C_WHITE),
        ]
        for i, (label, c) in enumerate(pipeline):
            x = 0.2 + i * 1.6
            add_rect(slide, x, 1.0, 1.45, 0.85, C_DARK_BOX)
            add_text(slide, label, x + 0.05, 1.1, 1.35, 0.65,
                     font_size=11, color=c, align=PP_ALIGN.CENTER)
            if i < 5:
                add_text(slide, "→", x + 1.47, 1.3, 0.15, 0.3,
                         font_size=12, bold=True, color=C_ACCENT,
                         align=PP_ALIGN.CENTER)

        # 검증 2가지
        add_rect(slide, 0.3, 2.1, 4.5, 2.5, C_DARK_BOX)
        add_text(slide, "이중 검증", 0.4, 2.15, 4.3, 0.38,
                 font_size=14, bold=True, color=C_GOLD)
        add_text(slide, "숫자 검증", 0.4, 2.6, 1.5, 0.35,
                 font_size=13, bold=True, color=C_RED)
        add_text(slide,
                 "LLM 응답의 수치가 실제 RSI14/이격도20/\n확률과 ±1.0 이내인지 확인",
                 0.4, 2.98, 4.2, 0.65, font_size=11, color=C_WHITE)
        add_text(slide, "인용 검증", 0.4, 3.7, 1.5, 0.35,
                 font_size=13, bold=True, color=C_RED)
        add_text(slide,
                 "(종목코드, 날짜) 형태 인용이 실제\nRAG 검색결과에 존재하는지 확인",
                 0.4, 4.08, 4.2, 0.55, font_size=11, color=C_WHITE)

        # 폴백
        add_rect(slide, 5.0, 2.1, 4.7, 2.5, C_DARK_BOX)
        add_text(slide, "Fallback 템플릿", 5.1, 2.15, 4.5, 0.38,
                 font_size=14, bold=True, color=C_GREEN)
        fb_cases = [
            ("API 오류",         "ANTHROPIC_API_KEY 없음"),
            ("숫자 불일치",      "검증 실패 시 자동 전환"),
            ("인용 오류",        "없는 사례 인용 감지"),
            ("폴백 출력",        "수치 기반 4문장 한국어 설명"),
        ]
        for i, (k, v) in enumerate(fb_cases):
            add_text(slide, k, 5.1, 2.62 + i * 0.48, 1.8, 0.4,
                     font_size=12, bold=True, color=C_ACCENT2)
            add_text(slide, v, 7.0, 2.62 + i * 0.48, 2.6, 0.4,
                     font_size=11, color=C_GRAY)

        # 시스템 프롬프트 핵심
        add_rect(slide, 0.3, 4.75, 9.4, 1.5, RGBColor(0x06, 0x14, 0x22))
        add_text(slide, "시스템 프롬프트 핵심 원칙", 0.5, 4.8, 4.0, 0.38,
                 font_size=13, bold=True, color=C_ACCENT)
        sp_lines = (
            "• 주어진 수치 데이터와 검색된 과거 사례만으로 설명 — 추측·새로운 판단 금지\n"
            "• 사례 인용 시 반드시 (종목코드, 날짜) 출처 표기\n"
            "• 검색 사례 없으면 '참고 사례를 찾지 못했습니다' 명시, 절대 지어내지 않음\n"
            "• 확신도 낮음 플래그 켜진 경우 반드시 언급"
        )
        add_text(slide, sp_lines, 0.5, 5.22, 9.1, 1.0,
                 font_size=11, color=C_WHITE)

    slide_content(prs, "STEP 8  LLM 자연어 설명 생성 (llm_explainer.py)", s10)

    # ── 11. Streamlit 챗봇 ────────────────────────────────────────────────────
    def s11(slide):
        # 대화 흐름
        add_rect(slide, 0.3, 0.95, 9.4, 0.55, C_DARK_BOX)
        flow = "사용자 입력 → 종목 추출 → PPO 추론 → 수치 근거 → RAG 검색 → LLM 설명 → 채팅 출력"
        add_text(slide, flow, 0.5, 1.02, 9.0, 0.4,
                 font_size=12, color=C_ACCENT, align=PP_ALIGN.CENTER)

        # UI 구성
        add_rect(slide, 0.3, 1.65, 4.5, 3.8, C_DARK_BOX)
        add_text(slide, "UI 구성", 0.4, 1.7, 4.3, 0.38,
                 font_size=14, bold=True, color=C_GOLD)
        ui_items = [
            ("사이드바",     "종목 selectbox\n모델 로딩 상태\n대화 초기화 버튼"),
            ("채팅창",       "매수/유보/매도 색상 배지\naction_probs 막대 표시"),
            ("Expander",     "Top-K 사례 원문\n사후 수익률 표시"),
            ("모델 부재 시", "입력 블록 + 안내 메시지"),
        ]
        for i, (k, v) in enumerate(ui_items):
            y = 2.18 + i * 0.82
            add_text(slide, k, 0.5, y, 1.5, 0.72,
                     font_size=12, bold=True, color=C_ACCENT2)
            add_text(slide, v, 2.1, y, 2.6, 0.72,
                     font_size=11, color=C_GRAY)

        # 출력 예시
        add_rect(slide, 5.0, 1.65, 4.7, 3.8, RGBColor(0x06, 0x14, 0x22))
        add_text(slide, "출력 예시", 5.1, 1.7, 4.5, 0.38,
                 font_size=14, bold=True, color=C_GREEN)
        example2 = (
            "사용자: 하이트진로 지금 사도 될까?\n\n"
            "[매수]  매수 62% / 유보 31% / 매도 7%\n\n"
            "현재 RSI14는 58.2로 중립 구간이며,\n"
            "이격도20도 102.1로 과열되지 않은 상태입니다.\n"
            "4일 전 골든크로스 발생 — 추세 추종 보너스 조건.\n\n"
            "과거 유사 사례(000080, 2024-11-12)에서도\n"
            "비슷한 패턴 이후 5일 +2.1% 상승 경향.\n\n"
            "▸ 참고한 과거 사례 (펼치기)"
        )
        add_text(slide, example2, 5.1, 2.15, 4.5, 3.2,
                 font_size=11, color=C_ACCENT2)

        # 면책 고지
        add_rect(slide, 0.3, 5.6, 9.4, 0.65, RGBColor(0x3B, 0x10, 0x10))
        add_text(slide,
                 "이 답변은 과거 데이터 기반 모델의 참고 지표이며 투자 자문이 아닙니다.\n"
                 "실제 매매는 본인 판단과 책임 하에 결정하세요.",
                 0.5, 5.67, 9.0, 0.55, font_size=12, color=C_RED,
                 align=PP_ALIGN.CENTER)

        add_text(slide, "streamlit run app.py", 3.5, 6.35, 3.0, 0.38,
                 font_size=12, bold=True, color=C_GRAY, align=PP_ALIGN.CENTER)

    slide_content(prs, "STEP 9  Streamlit 챗봇 완성 (app.py)", s11)

    # ── 12. 전체 데이터 흐름 요약 ────────────────────────────────────────────
    def s12(slide):
        # 주 트랙
        main_flow = [
            ("stock_prices.csv", 0.3),
            ("indicators/\ntechnical.py", 0.3),
            ("env/\ntrading_env.py", 0.3),
            ("train/\nrun_training.py", 0.3),
            ("models/\nbest_model.zip", 0.3),
            ("inference/\npredict.py", 0.3),
            ("explain/\nrule_based.py", 0.3),
        ]
        for i, (label, alpha) in enumerate(main_flow):
            x = 0.25 + i * 1.35
            add_rect(slide, x, 1.0, 1.2, 1.1, C_DARK_BOX)
            add_text(slide, label, x + 0.05, 1.1, 1.1, 0.9,
                     font_size=9, color=C_ACCENT, align=PP_ALIGN.CENTER)
            if i < 6:
                add_text(slide, "→", x + 1.22, 1.4, 0.15, 0.35,
                         font_size=10, bold=True, color=C_ACCENT, align=PP_ALIGN.CENTER)

        add_text(slide, "주(主) PPO 트랙", 0.3, 0.8, 5.0, 0.25,
                 font_size=12, bold=True, color=C_ACCENT)

        # 부 트랙
        rag_flow = [
            ("corpus/\nbuild_corpus.py", C_GREEN),
            ("chroma_db/\nbm25_index", C_GREEN),
            ("rag_retriever.py\n(4단계)", C_GREEN),
            ("llm_explainer.py\nClaude Haiku", C_GREEN),
        ]
        for i, (label, c) in enumerate(rag_flow):
            x = 0.25 + i * 2.4
            add_rect(slide, x, 2.5, 2.1, 1.1, C_DARK_BOX)
            add_text(slide, label, x + 0.05, 2.6, 2.0, 0.9,
                     font_size=10, color=c, align=PP_ALIGN.CENTER)
            if i < 3:
                add_text(slide, "→", x + 2.12, 2.9, 0.3, 0.35,
                         font_size=12, bold=True, color=C_GREEN, align=PP_ALIGN.CENTER)

        add_text(slide, "부(副) Advanced RAG 트랙", 0.3, 2.3, 5.0, 0.25,
                 font_size=12, bold=True, color=C_GREEN)

        # 연결 화살표
        add_text(slide, "↓  PPO 결정 후 호출", 7.5, 1.8, 2.5, 0.4,
                 font_size=11, bold=True, color=C_GOLD)

        # 최종 출력
        add_rect(slide, 0.3, 3.85, 9.4, 1.1, C_DARK_BOX)
        add_rect(slide, 0.3, 3.85, 0.08, 1.1, C_GOLD)
        add_text(slide, "결합 →", 0.5, 3.92, 1.0, 0.4,
                 font_size=13, bold=True, color=C_GOLD)
        add_text(slide,
                 "app.py (Streamlit 챗봇)  —  '[종목] 매수/매도/유보 + 과거 유사 사례 근거 + 자연어 설명'",
                 1.6, 3.92, 7.9, 0.45, font_size=13, color=C_WHITE)
        add_text(slide,
                 "PPO가 결정 → RAG가 설명 보강 → LLM이 자연어로 변환  |  "
                 "RAG는 PPO 결정을 바꾸지 않는다",
                 0.5, 4.38, 9.1, 0.5, font_size=11, color=C_GRAY)

        # 핵심 설계 원칙 요약
        principles = [
            ("srtnCd 경계 준수",   "지표·episode 계산에서 종목 간 데이터 혼합 절대 금지"),
            ("PPO 主 / RAG 副",   "RAG는 설명만, PPO 결정을 바꾸거나 선행 호출 금지"),
            ("결정론적 수치 근거", "rule_based.py — LLM 미사용, 데이터에서 직접 계산"),
            ("폴백 보장",          "LLM/RAG 실패 시 수치 기반 템플릿으로 자동 전환"),
        ]
        for i, (k, v) in enumerate(principles):
            col = i % 2
            row = i // 2
            x = 0.3 + col * 4.9
            y = 5.15 + row * 0.58
            add_rect(slide, x, y, 4.6, 0.5, RGBColor(0x0A, 0x20, 0x32))
            add_text(slide, k, x + 0.1, y + 0.06, 1.6, 0.38,
                     font_size=11, bold=True, color=C_GOLD)
            add_text(slide, v, x + 1.75, y + 0.06, 2.75, 0.38,
                     font_size=10, color=C_WHITE)

    slide_content(prs, "전체 데이터 흐름 요약 — 설계 핵심 원칙", s12)

    # ── 13. 마무리 슬라이드 ───────────────────────────────────────────────────
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_bg(slide, prs, C_BG)
    add_rect(slide, 0, 0, 10, 0.1, C_ACCENT)
    add_rect(slide, 0, 7.4, 10, 0.1, C_ACCENT)
    add_rect(slide, 0, 0, 0.12, 7.5, C_ACCENT)

    add_text(slide, "핵심 한 줄 요약", 0.5, 1.5, 9, 0.5,
             font_size=22, bold=True, color=C_ACCENT2)
    add_rect(slide, 0.5, 2.1, 9.0, 1.4, C_DARK_BOX)
    add_text(slide,
             "PPO가 먼저 매수/매도/유보를 결정하고,\n"
             "Advanced RAG가 과거 유사 사례로 그 이유를 설명하며,\n"
             "Claude가 자연어로 번역한다.",
             0.7, 2.2, 8.6, 1.2,
             font_size=16, color=C_WHITE, align=PP_ALIGN.CENTER)

    checklist = [
        ("완료", "데이터 수집 · 전처리 (srtnCd groupby)", C_GREEN),
        ("완료", "기술 지표 계산 (RSI · 이격도 · 크로스)", C_GREEN),
        ("완료", "TradingEnv (포지션 · 미실현손익 · 거래비용)", C_GREEN),
        ("완료", "PPO 학습 (200종목 단일 정책망)", C_GREEN),
        ("완료", "추론 + 수치 근거 (rule_based, 결정론적)", C_GREEN),
        ("완료", "RAG Corpus 구축 (ChromaDB + BM25)", C_GREEN),
        ("완료", "Advanced RAG 4단계 (Expansion·Hybrid·Filter·Rerank)", C_GREEN),
        ("완료", "LLM 설명 생성 (Claude Haiku + 이중 검증 + 폴백)", C_GREEN),
        ("완료", "Streamlit 챗봇 (종목명 입력 → 판단 + 근거)", C_GREEN),
    ]
    for i, (status, item, c) in enumerate(checklist):
        col = i % 2
        row = i // 2
        x = 0.5 + col * 4.85
        y = 3.7 + row * 0.55
        if i == 8:
            x = 0.5
        add_text(slide, f"✓  {item}", x, y, 4.6, 0.48,
                 font_size=11, color=c)

    add_text(slide, "금융 LLM 실습  Day 3  |  01_golden_cross",
             0.5, 7.05, 9, 0.35, font_size=12, color=C_GRAY,
             align=PP_ALIGN.CENTER)

    # ── 저장 ─────────────────────────────────────────────────────────────────
    out = "/home/hpe/finance-llm-course/day3/01_golden_cross/ppo_rag_system_design.pptx"
    prs.save(out)
    print(f"저장 완료: {out}")


if __name__ == "__main__":
    make_ppt()
